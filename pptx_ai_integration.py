# app.py - Main Flask Server (with PDF support)

#if using ngrok, change address in CORS to your ngrok address

from flask import Flask, request, render_template_string, jsonify
from pptx import Presentation
from flask_cors import CORS
import fitz  # PyMuPDF
import os
import requests


app = Flask(__name__)

#The third address has to change if your ngrok address changes
CORS(app,origins=["https://localhost:3000", "http://localhost:3000", os.getenv("REACT_APP_NGROK_URL")])
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)



#set to completions for now can change 
AI_SERVICE_URL = "https://api.openai.com/v1/chat/completions"  # Replace with your AI service endpoint

#set to this if you want an image
#AI_SERVICE_URL = "https://api.openai.com/v1/images/generations"

UPLOAD_FORM_HTML = """
<h2>Upload PowerPoint or PDF File</h2>
<form method="post" enctype="multipart/form-data">
  <input type="file" name="uploaded_file">
  <input type="submit" value="Upload">
</form>
"""

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        if 'uploaded_file' not in request.files:
            return jsonify({"error": "No file part"})

        file = request.files['uploaded_file']
        if file.filename == '':
            return jsonify({"error": "No selected file"})

        filename = file.filename.lower()
        if filename.endswith(('.pptx', '.pdf')):
            filepath = os.path.join(UPLOAD_FOLDER, file.filename)
            file.save(filepath)

            if filename.endswith('.pptx'):
                extracted_data = extract_pptx_data(filepath)
            elif filename.endswith('.pdf'):
                extracted_data = extract_pdf_data(filepath)

            os.remove(filepath)

            prompt=request.form.get('prompt', 'Summarize this presentation')
            if not prompt:
                prompt = "Summarize this presentation"

            ai_response = send_to_ai(extracted_data, prompt)
            print(ai_response)

            return jsonify({
                "status": "File processed and sent to AI",
                "ai_response": ai_response
            })

        return jsonify({"error": "Invalid file type. Please upload a .pptx or .pdf file."})
    

    return render_template_string(UPLOAD_FORM_HTML)

def extract_pptx_data(filepath):
    prs = Presentation(filepath)
    slide_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        content = {"slide_number": idx, "texts": []}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    content["texts"].append(text)
        if content["texts"]:
            slide_data.append(content)

    return slide_data

def extract_pdf_data(filepath):
    doc = fitz.open(filepath)
    page_data = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if text:
            page_data.append({
                "page_number": page_num,
                "texts": [line.strip() for line in text.split('\n') if line.strip()]
            })

    return page_data

# def send_to_ai(extracted_data):
#     try:
#         response = requests.post(
#             AI_SERVICE_URL,
#             json={"presentation_data": extracted_data},
#             timeout=15
#         )
#         return response.json() if response.ok else {"error": response.text}
#     except requests.exceptions.RequestException as e:
#         return {"error": str(e)}

# if __name__ == '__main__':
#     app.run(debug=True)

def send_to_ai(extracted_data, prompt):
    try:
        headers = {
            "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}",
            "Content-Type": "application/json"
        }
        full_prompt = f"{prompt}"+ str(extracted_data)
        data = {
            "model": "gpt-4o",
            "messages": [{"role": "user", "content": f"{full_prompt}"}],
            "max_tokens": 300
        }
        response = requests.post(
            AI_SERVICE_URL,
            headers=headers,
            json=data,
            timeout=15
        )
        return response.json() if response.ok else {"error": response.text}
    except requests.exceptions.RequestException as e:
        return {"error": str(e)}







#@app.route('/analyze', methods=['POST'])
def analyze():
    data = request.json
    presentation = data.get("presentation_data", [])
    visualization_results = run_ai_on(presentation)
    return jsonify({"visualization": visualization_results})

def run_ai_on(presentation):
    return f"Processed {len(presentation)} slides for visualization"


if __name__ == '__main__':
    app.run(debug=True, ssl_context=('cert.pem', 'key.pem'))


